from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation('Geekstorm_Presentation.pptx')

# Slide 6: System Workflow
s6 = prs.slides[5]
# Inject a textbox for the workflow
txBox = s6.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
tf = txBox.text_frame
tf.text = "System Workflow:"

p1 = tf.add_paragraph()
p1.text = "1. INTAKE: User imports flexible tasks & daily active window."
p2 = tf.add_paragraph()
p2.text = "2. COMPUTATION: The engine identifies exact free time gaps."
p3 = tf.add_paragraph()
p3.text = "3. ALLOCATION: Greedy algorithms assign high-priority tasks first."
p4 = tf.add_paragraph()
p4.text = "4. EXECUTION: Strict Focus Lock engages during active timers."
p5 = tf.add_paragraph()
p5.text = "5. RECALIBRATION: Missed tasks automatically shift schedule."

for paragraph in tf.paragraphs:
    paragraph.font.size = Pt(22)

# Slide 7 (Index 6) is Instructions. We need to delete it.
if len(prs.slides) >= 7:
    rId = prs.slides._sldIdLst[6].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[6]

prs.save('Geekstorm_Presentation_Final.pptx')
print("Successfully generated Geekstorm_Presentation_Final.pptx")
