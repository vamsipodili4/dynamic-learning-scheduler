from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def replace_text_in_shape(shape, old_text, new_text):
    if not shape.has_text_frame:
        return
    if old_text in shape.text:
        # Preserve formatting where possible by replacing at the paragraph/run level
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)

prs = Presentation('PPT Template.pptx')

# Slide 1:
s1 = prs.slides[0]
for shape in s1.shapes:
    replace_text_in_shape(shape, "Problem Statement ID –", "Problem Statement ID – GEEK-909")
    replace_text_in_shape(shape, "Problem Statement Title-", "Problem Statement Title - Dynamic Learning Scheduler")
    replace_text_in_shape(shape, "Team Name (Registered on portal)", "Team Name - Tech Titans")

# Slide 2:
s2 = prs.slides[1]
for shape in s2.shapes:
    replace_text_in_shape(shape, "Describe your problem statement analysis", 
    "• Static schedules break when a single 1-hour task is missed, causing student anxiety.\n"
    "• Students struggle to identify 'free blocks' in chaotic routines.\n"
    "• Traditional planners are visually sterile and lack psychological engagement.")

# Slide 3:
s3 = prs.slides[2]
for shape in s3.shapes:
    replace_text_in_shape(shape, "Describe about the challenge you are trying to solve",
    "• Architecting an 'Adaptive Engine' that auto-shifts tasks based on real-time performance.\n"
    "• Eliminating digital distractions during deep work using a Strict Focus Lock (Tab blocking).\n"
    "• Combining a gorgeous Dual-Theme 3D Aesthetic with gamification.")

# Slide 4: Key features
s4 = prs.slides[3]
# The user screenshot just shows "Key features" at the top with an empty grey body.
# We will inject a textbox into the middle of the slide.
txBox = s4.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
tf = txBox.text_frame
tf.text = "1. Real-Time Adaptive Engine: Re-slots missed tasks automatically."
p2 = tf.add_paragraph()
p2.text = "2. Hybrid Scheduling: Users can manually 'Claim' free slots."
p3 = tf.add_paragraph()
p3.text = "3. Strict Focus Trap: Page Visibility API detects distractions."
p4 = tf.add_paragraph()
p4.text = "4. Dual-Theme UX: Sunlit Dashboard and Premium Spline Landing Page."

# Make the font format look okay
for paragraph in tf.paragraphs:
    paragraph.font.size = Pt(20)

# Slide 5: Prototype Link
s5 = prs.slides[4]
# Assuming we can just append a text box if "Proposed Prototype link" is just the title
txBox2 = s5.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
tf2 = txBox2.text_frame
tf2.text = "• Live Server: http://localhost:5173"
p = tf2.add_paragraph()
p.text = "• Repository: https://github.com/vamsipodili/Dynamic-Learning-Scheduler"
for paragraph in tf2.paragraphs:
    paragraph.font.size = Pt(24)

prs.save('Geekstorm_Presentation.pptx')
print("Successfully generated Geekstorm_Presentation.pptx")
